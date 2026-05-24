#!/usr/bin/env python3
# orchestrator.py — NotebookLM-style PDF Pipeline
#
# Flow:
#   Any input (.txt / .csv / .pdf)
#     → [Agent 0+1: Analyzer + Planner]  Understand content, plan N slides
#     → [Agent 2:   Designer]            Generate HTML per slide (with inline SVG icons)
#     → [Agent 3:   Assembler]           Merge into one print-ready HTML
#     → [Agent 4:   Critic]              Score → feedback
#     → if score ≥ threshold → export PDF
#     → if score <  threshold → patch broken slides → loop
#
# Usage:
#   python orchestrator.py --input alerts.txt
#   python orchestrator.py --input topic.txt --output output/report.pdf
#   python orchestrator.py --input data.csv --style dark --html-only
#   python orchestrator.py --input logs.txt --iterations 5 --threshold 8.0

import argparse, json, os, sys
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from pathlib import Path

try:
    from rich.console import Console
    from rich.rule    import Rule
    from rich.table   import Table
    from rich         import box
    C = Console()
    def log(msg, style=""): C.print(msg, style=style)
    def rule(t=""): C.print(Rule(t))
    HAS_RICH = True
except ImportError:
    HAS_RICH = False
    def log(msg, style=""): print(msg)
    def rule(t=""): print(f"\n{'─'*60} {t}")

import config
sys.path.insert(0, str(Path(__file__).parent))


class PipelineInputError(RuntimeError):
    """Invalid or missing input (API-safe; CLI maps to exit code 1)."""


class PipelinePlanError(RuntimeError):
    """Planner produced no slides."""


# ══════════════════════════════════════════════════════════════════
#  INPUT LOADER
# ══════════════════════════════════════════════════════════════════

def load(path: str) -> str:
    p = Path(path)
    if not p.exists():
        log(f"[red]File not found: {path}[/red]")
        raise PipelineInputError(f"File not found: {path}")

    ext = p.suffix.lower()
    log(f"📂 Loading [cyan]{ext}[/cyan] — {p.name}")

    # Image — Gemma 4 multimodal: extract all text/data from the image via vision API
    if ext in ('.png', '.jpg', '.jpeg', '.webp', '.gif'):
        import base64
        with open(path, 'rb') as f:
            img_bytes = f.read()
        _MIME = {'.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                 '.png': 'image/png', '.webp': 'image/webp', '.gif': 'image/gif'}
        mime = _MIME.get(ext, 'image/jpeg')
        log(f"   {len(img_bytes):,} bytes — extracting via Gemma 4 vision")
        try:
            from utils.llm import call_vision
            text = call_vision(
                "You are a data extractor. Extract ALL text, numbers, labels, table data, "
                "chart values, and any other information visible in this image. "
                "Preserve structure and formatting. Be exhaustive — miss nothing.",
                base64.b64encode(img_bytes).decode(),
                mime,
                key="analyzer",
                max_tokens=4000,
            )
            log(f"   Vision extracted {len(text):,} chars")
            return text[:config.MAX_DATA_CHARS]
        except Exception as e:
            log(f"   [yellow]Vision extraction failed: {e}[/yellow]")
            return f"[Image file: {p.name}. Content could not be extracted: {e}]"

    # PDF
    if ext == '.pdf':
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)
            log(f"   {len(text):,} chars from PDF ({len(pdf.pages)} pages)")
            return text[:config.MAX_DATA_CHARS]
        except ImportError:
            pass
        try:
            import PyPDF2
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = "\n".join(
                    reader.pages[i].extract_text() or ""
                    for i in range(len(reader.pages))
                )
            return text[:config.MAX_DATA_CHARS]
        except ImportError:
            log("   [red]No PDF library. Install: pip install pdfplumber[/red]")
            raise PipelineInputError(
                "No PDF library. Install: pip install pdfplumber (or PyPDF2)"
            ) from None

    # CSV
    if ext == '.csv':
        try:
            import pandas as pd
            df  = pd.read_csv(path)
            raw = f"CSV data — {len(df)} rows × {len(df.columns)} columns\nColumns: {list(df.columns)}\n\n"
            raw += df.to_string(index=False)
            log(f"   {len(df)} rows × {len(df.columns)} cols")
            return raw[:config.MAX_DATA_CHARS]
        except ImportError:
            pass  # Fall through to text reader

    # Default: plain text
    with open(path, encoding='utf-8', errors='replace') as f:
        raw = f.read()
    truncated = len(raw) > config.MAX_DATA_CHARS
    log(f"   {len(raw):,} chars{' (truncating)' if truncated else ''}")
    return raw[:config.MAX_DATA_CHARS] + ("\n...[truncated]" if truncated else "")


# ══════════════════════════════════════════════════════════════════
#  PDF EXPORT
# ══════════════════════════════════════════════════════════════════

def _playwright_html_to_pdf(abs_html: str, pdf_path: str) -> None:
    """Sync Playwright run — must execute in a worker thread under Uvicorn/asyncio."""
    from playwright.sync_api import sync_playwright

    log("   Converting via [cyan]Playwright[/cyan]...")
    with sync_playwright() as pw:
        br = pw.chromium.launch()
        try:
            pg = br.new_page(viewport={"width": 1280, "height": 720})
            pg.goto(f"file:///{abs_html}", wait_until="networkidle")

            chart_total = pg.evaluate(
                "parseInt(document.querySelector('meta[name=chart-total]')?.content || '0')"
            )
            if chart_total > 0:
                log(f"   Waiting for {chart_total} Chart.js charts...")
                try:
                    pg.wait_for_function(
                        f"window.__chartsReady >= {chart_total}", timeout=10000
                    )
                except Exception:
                    log("   [yellow]Chart wait timed out — capturing anyway[/yellow]")
            else:
                pg.wait_for_timeout(1500)

            pg.wait_for_timeout(500)

            pg.pdf(
                path=pdf_path,
                width="1280px",
                height="720px",
                print_background=True,
                margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
            )
        finally:
            br.close()


def _playwright_html_to_pptx(abs_html: str, pptx_path: str) -> None:
    """Screenshot each slide via Playwright, build PPTX with python-pptx."""
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Emu
    import io

    SLIDE_W, SLIDE_H = 1280, 720
    log("   Screenshotting slides via [cyan]Playwright[/cyan]...")
    with sync_playwright() as pw:
        br = pw.chromium.launch()
        try:
            # device_scale_factor=1 ensures screenshots are exactly SLIDE_W×SLIDE_H
            # pixels regardless of the host display's DPI.
            ctx = br.new_context(
                viewport={"width": SLIDE_W, "height": SLIDE_H},
                device_scale_factor=1.0,
            )
            pg = ctx.new_page()
            pg.goto(f"file:///{abs_html}", wait_until="networkidle")
            chart_total = pg.evaluate(
                "parseInt(document.querySelector('meta[name=chart-total]')?.content || '0')"
            )
            if chart_total > 0:
                log(f"   Waiting for {chart_total} Chart.js charts...")
                try:
                    pg.wait_for_function(
                        f"window.__chartsReady >= {chart_total}", timeout=10000
                    )
                except Exception:
                    log("   [yellow]Chart wait timed out — capturing anyway[/yellow]")
            else:
                pg.wait_for_timeout(1500)
            pg.wait_for_timeout(500)
            slide_els = pg.query_selector_all("section.slide")
            screenshots = [el.screenshot() for el in slide_els]
        finally:
            br.close()

    prs = Presentation()
    prs.slide_width  = Emu(12192000)  # 13.333" at 96 dpi = 1280 px
    prs.slide_height = Emu(6858000)   # 7.5"   at 96 dpi = 720 px
    blank_layout = prs.slide_layouts[6]
    for png_bytes in screenshots:
        sl = prs.slides.add_slide(blank_layout)
        sl.shapes.add_picture(io.BytesIO(png_bytes), 0, 0, prs.slide_width, prs.slide_height)
    prs.save(pptx_path)


def export_pptx(html_path: str, pptx_path: str) -> bool:
    abs_html = str(Path(html_path).resolve())
    try:
        import importlib.util
        if importlib.util.find_spec("playwright") is None:
            raise ImportError("playwright not installed")
        if importlib.util.find_spec("pptx") is None:
            raise ImportError("python-pptx not installed")
        with ThreadPoolExecutor(max_workers=1, thread_name_prefix="playwright_pptx") as pool:
            pool.submit(_playwright_html_to_pptx, abs_html, pptx_path).result(timeout=600)
        log(f"   [green]✅ PPTX: {pptx_path}[/green]")
        return True
    except ImportError as e:
        log(f"   [yellow]PPTX export unavailable: {e}[/yellow]")
    except Exception as e:
        log(f"   [yellow]PPTX export failed: {e}[/yellow]")
    return False


def export_pdf(html_path: str, pdf_path: str) -> bool:
    abs_html = str(Path(html_path).resolve())

    # 1. Playwright (best — waits for Chart.js renders)
    # Run sync API in a dedicated thread: Uvicorn runs handlers on asyncio; Playwright
    # sync API must not run on the event-loop thread ("Sync API inside the asyncio loop").
    try:
        import importlib.util

        if importlib.util.find_spec("playwright") is None:
            raise ImportError("playwright not installed")
        with ThreadPoolExecutor(max_workers=1, thread_name_prefix="playwright_pdf") as pool:
            pool.submit(_playwright_html_to_pdf, abs_html, pdf_path).result(timeout=600)
        log(f"   [green]✅ PDF: {pdf_path}[/green]")
        return True
    except ImportError:
        pass
    except Exception as e:
        log(f"   [yellow]Playwright failed: {e}[/yellow]")

    # 2. pdfkit fallback
    try:
        import pdfkit
        opts = {
            "page-width": "1280px", "page-height": "720px",
            "enable-local-file-access": "",
            "javascript-delay": "4000",
            "quiet": "",
        }
        log("   Converting via [cyan]pdfkit[/cyan]...")
        pdfkit.from_file(abs_html, pdf_path, options=opts)
        log(f"   [green]✅ PDF: {pdf_path}[/green]")
        return True
    except ImportError:
        pass
    except Exception as e:
        log(f"   [yellow]pdfkit failed: {e}[/yellow]")

    log(f"""
  [yellow]⚠ No PDF converter available.[/yellow]
  Install Playwright:   pip install playwright && playwright install chromium

  Or open in Chrome → File → Print → Save as PDF:
    [cyan]{abs_html}[/cyan]
""")
    return False


# ══════════════════════════════════════════════════════════════════
#  PPTX THEME UNIFICATION
# ══════════════════════════════════════════════════════════════════

# Higher priority = preferred for a cohesive professional presentation
_PPTX_PERSONA_PRIORITY: dict[str, int] = {
    "editorial_dark":    8,
    "technical_dense":   7,
    "data_dashboard":    6,
    "minimalist_focus":  5,
    "narrative_warm":    4,
    "magazine_spread":   3,
    "vibrant_split":     2,
    "infographic_vibrant": 1,
}


def _unify_pptx_personas(plan: dict) -> None:
    """Collapse per-slide aesthetic_persona to one consistent theme for PPTX.

    PDF variety (8 personas) looks intentional; PPTX slide-flipping makes
    background jumps look broken.  Pick the single most-professional persona
    from those already assigned by the planner and stamp it on every slide.
    """
    slides = plan.get("slides", [])
    if not slides:
        return
    from collections import Counter
    counts: Counter = Counter(
        s.get("aesthetic_persona") for s in slides if s.get("aesthetic_persona")
    )
    if not counts:
        return
    # Sort by (frequency desc, priority desc) so we get the most-used
    # professional persona in case of a tie.
    theme = max(
        counts,
        key=lambda p: (counts[p], _PPTX_PERSONA_PRIORITY.get(p, 0)),
    )
    log(
        f"  [PPTX] Unifying theme → [cyan]{theme}[/cyan] "
        f"(was {len(counts)} different personas across {len(slides)} slides)"
    )
    for slide in slides:
        slide["aesthetic_persona"] = theme


# ══════════════════════════════════════════════════════════════════
#  MAIN PIPELINE
# ══════════════════════════════════════════════════════════════════

def run(input_path: str, output_path: str, html_only: bool = False, output_format: str = "pdf"):
    out_dir = Path(output_path).parent
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = Path(output_path).stem
    ts   = datetime.now().strftime("%H:%M:%S")

    rule(f"🚀  NotebookLM PDF Pipeline  [{ts}]")
    log(f"  Input:      [cyan]{input_path}[/cyan]")
    log(f"  Output:     [cyan]{output_path}[/cyan]")
    log(f"  Style:      [cyan]{config.VISUAL_STYLE}[/cyan]")
    log(f"  Seed:       [cyan]{config.DESIGN_SEED}[/cyan]  (changes every run)")
    log(f"  Max loops:  [cyan]{config.MAX_ITERATIONS}[/cyan]")
    log(f"  Threshold:  [cyan]{config.PASS_THRESHOLD}/10[/cyan]")
    rule()

    from agents.planner   import run as plan_slides
    from agents.designer  import run as design_slides
    from agents.assembler import run as assemble_html
    from agents.critic    import run as critique
    from agents.browser   import run as browser_enrich

    # ── AGENT 0.5: Browser / Web Research (optional) ──────────────
    raw_data = load(input_path)
    if config.BROWSER_ENABLED:
        rule("Agent 0.5 — Browser / Web Research")
        raw_data = browser_enrich(raw_data)
    else:
        log("  [Browser] Disabled — set BROWSER_ENABLED=true in .env to enable web research")

    # ── AGENTS 0+1: Analyzer + Planner ────────────────────────────
    rule("Agent 0+1 — Analyzer + Planner")
    plan = plan_slides(raw_data)

    if not plan.get('slides'):
        log("[red]Planner returned no slides — aborting[/red]")
        raise PipelinePlanError("Planner returned no slides")

    # For PPTX: collapse all persona variety to one cohesive theme.
    # PDF benefits from variety; PPTX slide-flipping makes background
    # jumps look broken.
    if output_format == "pptx":
        _unify_pptx_personas(plan)

    # Save plan
    plan_path = str(out_dir / f"{stem}_plan.json")
    plan_save = {k: v for k, v in plan.items() if not k.startswith('_')}
    with open(plan_path, 'w') as f:
        json.dump(plan_save, f, indent=2, default=str)
    log(f"  💾 Plan: [dim]{plan_path}[/dim]")

    # ── AGENTS 2+3+4: Design → Assemble → Critic loop ─────────────
    rule("Agent 2 — Designer  →  Agent 3 — Assembler  →  Agent 4 — Critic")

    best_html     = None
    best_score    = 0.0
    best_sections = None
    prev_sections = None
    history       = []
    feedback      = None

    for iteration in range(1, config.MAX_ITERATIONS + 1):
        log(f"\n  [bold]── Iteration {iteration}/{config.MAX_ITERATIONS} ──[/bold]")

        # Determine which slides to redo
        if iteration == 1:
            log("  [Agent 2] Designing all slides...")
            slides_to_redo = None
        else:
            fixes      = feedback.get('slides_to_fix', []) if feedback else []
            fix_slots  = [int(s.get('slot', 0)) for s in fixes if s.get('slot')]
            if fix_slots:
                log(f"  [Agent 2] Patching slots: {fix_slots} ({len(fix_slots)} of {len(plan['slides'])})")
                slides_to_redo = fix_slots
            else:
                log("  [Agent 2] No specific slots flagged — redesigning all")
                slides_to_redo = None

        # ── Agent 2: Designer ──────────────────────────────────────
        slides_html = design_slides(
            plan           = plan,
            feedback       = feedback,
            slides_to_redo = slides_to_redo,
            output_format  = output_format,
        )

        # ── Agent 3: Assembler ─────────────────────────────────────
        log("  [Agent 3] Assembling HTML...")
        html, sections = assemble_html(
            plan          = plan,
            slides_html   = slides_html,
            prev_sections = prev_sections,
        )
        prev_sections = sections

        # Save iteration HTML
        iter_html = str(out_dir / f"{stem}_iter{iteration}.html")
        with open(iter_html, 'w', encoding='utf-8') as f:
            f.write(html)
        log(f"  💾 [dim]{iter_html}[/dim]")

        # ── Agent 4: Critic ────────────────────────────────────────
        result = critique(html, plan, output_format=output_format)
        history.append(result.to_dict())

        if result.weighted_score > best_score:
            best_score    = result.weighted_score
            best_html     = html
            best_sections = sections

        crit_path = str(out_dir / f"{stem}_iter{iteration}_critic.json")
        with open(crit_path, 'w') as f:
            json.dump(result.to_dict(), f, indent=2)

        if result.passed:
            log(f"  [green]✅ Quality gate passed! {result.weighted_score:.2f} ≥ {config.PASS_THRESHOLD}[/green]")
            best_html = html
            break

        if iteration < config.MAX_ITERATIONS:
            n = len(result.slides_to_fix)
            log(
                f"  [yellow]Score {result.weighted_score:.2f} < {config.PASS_THRESHOLD} — "
                f"{'patching ' + str(n) + ' slide(s)' if n else 'full redesign'} next[/yellow]"
            )
            feedback = result.to_dict()
        else:
            log(f"  [yellow]Max iterations reached. Best score: {best_score:.2f}[/yellow]")

    # ── Export ─────────────────────────────────────────────────────
    rule("📄  Exporting")
    final_html = str(out_dir / f"{stem}.html")
    with open(final_html, 'w', encoding='utf-8') as f:
        f.write(best_html)
    log(f"  ✅ HTML: [green]{final_html}[/green]")

    if not html_only:
        if output_format == "pptx":
            export_pptx(final_html, output_path)
        elif output_path.endswith('.pdf'):
            export_pdf(final_html, output_path)

    # ── Summary ────────────────────────────────────────────────────
    rule("📊  Summary")
    if HAS_RICH and history:
        t = Table(box=box.SIMPLE, show_header=True)
        t.add_column("Iter", justify="center")
        t.add_column("Score", justify="center")
        t.add_column("Status")
        t.add_column("Verdict")
        for i, h in enumerate(history, 1):
            sc  = h["weighted_score"]
            col = "green" if h["passed"] else ("yellow" if sc >= 6 else "red")
            t.add_row(
                str(i),
                f"[{col}]{sc:.2f}[/{col}]",
                "✅ Pass" if h["passed"] else "↩ Retry",
                h["verdict"][:60],
            )
        C.print(t)

    log(f"  Final score: [bold]{best_score:.2f}/10[/bold]")
    log(f"  HTML:        [green]{final_html}[/green]")
    log(f"  Design seed: [dim]{config.DESIGN_SEED}[/dim] (run again for different layout)")


# ══════════════════════════════════════════════════════════════════
#  CLI
# ══════════════════════════════════════════════════════════════════

def main():
    p = argparse.ArgumentParser(description="NotebookLM-style PDF Generator")
    p.add_argument("--input",      required=True,  help="Input file (.txt/.csv/.pdf/.png/.jpg/.webp)")
    p.add_argument("--output",     default="output/report.pdf")
    p.add_argument("--iterations", type=int,   default=config.MAX_ITERATIONS)
    p.add_argument("--threshold",  type=float, default=config.PASS_THRESHOLD)
    p.add_argument("--html-only",  action="store_true")
    p.add_argument("--style",
                   default=config.VISUAL_STYLE,
                   choices=["notebooklm", "modern", "dark", "auto"],
                   help="Visual style preset (auto = tone-adaptive)")
    p.add_argument("--seed",       type=int, default=None,
                   help="Design seed (omit for random — same seed = same layout)")
    args = p.parse_args()

    config.MAX_ITERATIONS = args.iterations
    config.PASS_THRESHOLD = args.threshold
    config.VISUAL_STYLE   = args.style
    if args.seed is not None:
        config.DESIGN_SEED = args.seed

    try:
        run(args.input, args.output, html_only=args.html_only)
    except (PipelineInputError, PipelinePlanError) as e:
        log(f"[red]{e}[/red]")
        sys.exit(1)


if __name__ == "__main__":
    main()